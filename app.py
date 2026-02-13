# app.py
"""
Streamlit app: FC Den Bosch player PPTX/PDF generator

Repo structure (case-sensitive on Streamlit Cloud / Linux):
.
├─ app.py
├─ powerpoint_template.pptx
├─ bench.csv
├─ speler_foto's/
│   ├─ Kevin Monzialo.png
│   ├─ Kevin Felida.png
│   └─ ...
└─ requirements.txt

Streamlit Cloud secrets (FLAT keys; matches your current setup):
token_url = "..."
grant_type = "password"
username = "..."
password = "..."
client_id = "..."
client_secret = "..."
scope = ""
base_url = "https://..."   # API base used by api_get_json

Template placeholders supported:
- Text tokens: {{PLAYER_NAME}}, {{NATIONALITY}}, {{BIRTH_DATE}}, {{HEIGHT_M}}, {{CONTRACT}}, {{IS_EU}}, ...
- Image placeholders:
  - {IMAGE}                -> local player photo (speler_foto's/<name>.png), fallback to API imageUrl
  - {{RADAR_CHART}}        -> radar chart from bench.csv
  - {{PERFORMANCE_CHART}}  -> performance chart (uploaded OR auto-generated)
  - {PRESTATIES_FIGURE}    -> same as performance chart bytes (optional)

Notes:
- Token replacement handles "split across runs" (PowerPoint quirks).
- PDF export needs LibreOffice ("soffice") installed on the system PATH.
"""

from __future__ import annotations


import io
import math
import os
import re
import shutil
import subprocess
import tempfile
import unicodedata
from dataclasses import dataclass
from difflib import SequenceMatcher
from io import BytesIO
from typing import Any, Dict, List, Optional, Tuple

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import requests
import streamlit as st
from pptx import Presentation
from pathlib import Path


# ----------------------------
# Hardcoded FC Den Bosch players (for now)
# ----------------------------
FC_DEN_BOSCH_PLAYERS = [
    {"name": "Kevin Monzialo", "player_id": 40665},
    {"name": "Kevin Felida", "player_id": 35836},
    {"name": "Mees Laros", "player_id": 707170},
    {"name": "Ilias Boumassaoudi", "player_id": 678422},
    {"name": "Thijs van Leeuwen", "player_id": 52187},
    {"name": "Pepijn van de Merbel", "player_id": 52270},
    {"name": "Nick de Groot", "player_id": 65019},
    {"name": "Jeffry Fortes", "player_id": 8693},
    {"name": "Reda Akmum", "player_id": 46062},
    {"name": "Teun van Grunsven", "player_id": 45372},
    {"name": "Stan Maas", "player_id": 55461},
    {"name": "Sheddy Barglan", "player_id": 672112},
    {"name": "Genrich Sillé", "player_id": 708226},
    {"name": "Emian-Johar Semedo", "player_id": 713819},
    {"name": "Bohao Wang", "player_id": 697888},
    {"name": "Sebastian Karlsson Grach", "player_id": 58320},
    {"name": "Danny Verbeek", "player_id": 7301},
    {"name": "Jack de Vries", "player_id": 162192},
    {"name": "Zaid el Bakkali", "player_id": 719813},
    {"name": "Luc van Koeverden", "player_id": 673690},
]

CUSTOM_MAXES = {
    "Total distance (m)": 13750,
    "HI distance (m)": 1650,
    "Sprint distance (m)": 600,
    "HI runs": 50,
    "Sprint runs": 25,
}

SEASON_RE = re.compile(r"\b(20\d{2})\s*[/\-]\s*(20\d{2})\b")
def _safe_int(x: Any) -> int:
    try:
        if x is None:
            return 0
        return int(round(float(x)))
    except Exception:
        return 0


def _keys_preview(d: Any, max_keys: int = 25) -> str:
    if not isinstance(d, dict):
        return ""
    ks = list(d.keys())
    ks_sorted = sorted([str(k) for k in ks])[:max_keys]
    suffix = " ..." if len(ks) > max_keys else ""
    return ", ".join(ks_sorted) + suffix


def _get_team_name(it: dict) -> str:
    return ((it.get("team") or {}).get("name") or "").strip() or "Unknown"


def _get_comp_name(it: dict) -> str:
    # Competition naming differs per endpoint. This tries common patterns.
    comp = it.get("competition") or it.get("tournament") or it.get("league") or {}
    name = ""
    if isinstance(comp, dict):
        name = (comp.get("name") or "").strip()
    if not name:
        # last resort: scan nested strings for something that looks like a comp name
        for k in ["competitionName", "tournamentName", "leagueName", "name"]:
            v = it.get(k)
            if isinstance(v, str) and v.strip():
                name = v.strip()
                break
    return name or ""


def _pick_value_source(it: dict, stat_keys: List[str], metric_keys: List[str]) -> Tuple[int, str, str]:
    """
    Returns (value, source, key_used)
    source in {"stats","metrics","other","missing"}
    """
    stats = it.get("stats") or {}
    metrics = it.get("metrics") or {}

    # try stats
    if isinstance(stats, dict):
        for k in stat_keys:
            if k in stats and stats.get(k) is not None:
                return _safe_int(stats.get(k)), "stats", k

    # try metrics
    if isinstance(metrics, dict):
        for k in metric_keys:
            if k in metrics and metrics.get(k) is not None:
                return _safe_int(metrics.get(k)), "metrics", k

    # case-insensitive search (sometimes keys differ in casing)
    if isinstance(stats, dict):
        low = {str(k).lower(): k for k in stats.keys()}
        for k in stat_keys:
            lk = k.lower()
            if lk in low and stats.get(low[lk]) is not None:
                return _safe_int(stats.get(low[lk])), "stats", str(low[lk])

    if isinstance(metrics, dict):
        low = {str(k).lower(): k for k in metrics.keys()}
        for k in metric_keys:
            lk = k.lower()
            if lk in low and metrics.get(low[lk]) is not None:
                return _safe_int(metrics.get(low[lk])), "metrics", str(low[lk])

    # maybe on top-level
    for k in stat_keys + metric_keys:
        if k in it and it.get(k) is not None:
            return _safe_int(it.get(k)), "other", k

    return 0, "missing", ""


def debug_season_rows(api_base: str, token: str, player_id: int, season_id: int) -> None:
    """
    Prints a row-by-row inspection for BOTH endpoints you use:
      - contribution-ratings (games/minutes usually)
      - career-stats (goals/assists usually; sometimes games/minutes too)
    Shows where each value is coming from (stats vs metrics).
    """
    import pandas as pd
    import streamlit as st

    st.subheader(f"DEBUG: player={player_id} season_id={season_id}")

    endpoints = [
        ("contribution-ratings", "/api/v2/metrics/players/contribution-ratings"),
        ("career-stats", "/api/v2/metrics/career-stats/players"),
    ]

    for label, path in endpoints:
        st.markdown(f"### Endpoint: `{label}`")

        obj = api_get_json(
            api_base,
            token,
            path,
            params={"PlayerIds": player_id, "SeasonIds": [season_id], "Limit": 2000},
        )
        items = items_of(obj)

        rows = []
        for it in items:
            team = _get_team_name(it)
            comp = _get_comp_name(it)

            # Try to detect values and their sources
            games, games_src, games_key = _pick_value_source(
                it,
                stat_keys=["matchesPlayed", "matchPlayed", "matches", "games"],
                metric_keys=["matchesPlayed", "matchPlayed", "matches", "games"],
            )
            minutes, min_src, min_key = _pick_value_source(
                it,
                stat_keys=["minutesPlayed", "minutes"],
                metric_keys=["minutesPlayed", "minutes"],
            )
            goals, goal_src, goal_key = _pick_value_source(
                it,
                stat_keys=["goal", "goals", "goalNonPenalty", "goal_non_penalty"],
                metric_keys=["goal", "goals", "goalNonPenalty", "goal_non_penalty"],
            )
            assists, ast_src, ast_key = _pick_value_source(
                it,
                stat_keys=["assist", "assists"],
                metric_keys=["assist", "assists"],
            )

            rows.append(
                {
                    "team": team,
                    "competition": comp,
                    "games": games,
                    "games_src": games_src,
                    "games_key": games_key,
                    "minutes": minutes,
                    "minutes_src": min_src,
                    "minutes_key": min_key,
                    "goals": goals,
                    "goals_src": goal_src,
                    "goals_key": goal_key,
                    "assists": assists,
                    "assists_src": ast_src,
                    "assists_key": ast_key,
                    "stats_keys": _keys_preview(it.get("stats") or {}),
                    "metrics_keys": _keys_preview(it.get("metrics") or {}),
                }
            )

        df = pd.DataFrame(rows)
        if df.empty:
            st.warning("No items returned for this season.")
            continue

        # Helpful grouping to see if league rows exist but are zeros
        st.write("Rows (raw):")
        st.dataframe(df, use_container_width=True)

        st.write("Grouped sum by team + competition (so you can see which comps exist):")
        grouped = (
            df.groupby(["team", "competition"], dropna=False)[["games", "minutes", "goals", "assists"]]
            .sum()
            .reset_index()
            .sort_values(["team", "competition"])
        )
        st.dataframe(grouped, use_container_width=True)

        st.write("Count by source (where values are coming from):")
        src_counts = {
            "games": df["games_src"].value_counts().to_dict(),
            "minutes": df["minutes_src"].value_counts().to_dict(),
            "goals": df["goals_src"].value_counts().to_dict(),
            "assists": df["assists_src"].value_counts().to_dict(),
        }
        st.json(src_counts)

def _is_toto_knvb_beker_nld(item: dict) -> bool:
    """
    Returns True if an API item represents the Dutch cup: 'TOTO KNVB beker (NLD)'.
    Works across endpoints by checking common nested fields.
    """
    if not isinstance(item, dict):
        return False

    def _collect_strings(obj: Any, out: List[str]) -> None:
        if obj is None:
            return
        if isinstance(obj, str):
            out.append(obj)
            return
        if isinstance(obj, dict):
            for v in obj.values():
                _collect_strings(v, out)
            return
        if isinstance(obj, list):
            for v in obj:
                _collect_strings(v, out)

    strings: List[str] = []
    _collect_strings(item, strings)
    blob = " ".join(s.lower() for s in strings if isinstance(s, str))

    # Match both with/without parentheses, tolerate spacing/casing differences.
    has_name = ("toto knvb beker" in blob) or ("knvb beker" in blob)
    has_nld = ("(nld)" in blob) or (" nld" in blob) or ("countrycode nld" in blob) or ("alpha3code nld" in blob)
    return bool(has_name and has_nld)
  
def normalize_season_label(name: str) -> str:
    """
    Normalize season strings to 'YYYY/YYYY' when possible.
    Examples:
      '2023/2024 Regular Season' -> '2023/2024'
      '2023-2024' -> '2023/2024'
      '2023/24' -> '' (unknown)  # add more logic if needed
    """
    if not name:
        return ""
    m = SEASON_RE.search(str(name))
    if not m:
        return ""
    return f"{m.group(1)}/{m.group(2)}"

# ----------------------------
# App configuration (match your repo file names exactly)
# ----------------------------
TEMPLATE_PPTX_PATH = "powerpoint template.pptx"
BENCH_CSV_PATH = "df_bench.csv"
PLAYER_PHOTOS_DIR = "Speler_fotos"

TOKEN_HINT = "Generate access token first (required)."


# ----------------------------
# Radar configuration
# ----------------------------
RADAR_METRICS_MAP: Dict[str, str] = {
    "Total distance (m)": "total_distance",
    "HI distance (m)": "high_intensity_distance",
    "Sprint distance (m)": "sprint_distance",
    "HI runs": "hi_runs",
    "Sprint runs": "sprint_runs",
}

# CUSTOM_MAXES: Dict[str, float] = {
#     "Total distance (m)": 15000,
#     "HI distance (m)": 1500,
#     "Sprint distance (m)": 500,
#     "HI runs": 75,
#     "Sprint runs": 30,
# }

NAME_R_MULT = {
    "Total distance (m)": 1.22,
    "Sprint distance (m)": 1.22,
    "HI runs": 1.22,
}
DEFAULT_NAME_R = 1.33


@dataclass(frozen=True)
class SectionSecrets:
    token_url: str
    grant_type: str
    username: str
    password: str
    client_id: str
    client_secret: str
    scope: str
    base_url: str


def load_section_secrets() -> SectionSecrets:
    if "auth" not in st.secrets:
        raise RuntimeError('Missing [auth] in Streamlit Secrets.')

    auth = st.secrets["auth"]

    # allow both client_secret and the common typo client_secrete
    client_secret = auth.get("client_secret") or auth.get("client_secrete")

    required = ["token_url", "grant_type", "username", "password", "client_id", "base_url"]
    missing = [k for k in required if not auth.get(k)]
    if not client_secret:
        missing.append("client_secret (or client_secrete)")

    if missing:
        raise RuntimeError(f"Missing required auth keys: {missing}")

    return SectionSecrets(
        token_url=str(auth["token_url"]),
        grant_type=str(auth.get("grant_type", "password")),
        username=str(auth["username"]),
        password=str(auth["password"]),
        client_id=str(auth["client_id"]),
        client_secret=str(client_secret),
        scope=str(auth.get("scope", "")),
        base_url=str(auth["base_url"]),
    )


def generate_access_token_from_secrets() -> str:
    cfg = load_section_secrets()
    st.session_state["api_base"] = cfg.base_url
    data = {
        "grant_type": cfg.grant_type or "password",
        "username": cfg.username,
        "password": cfg.password,
        "client_id": cfg.client_id,
        "client_secret": cfg.client_secret,
    }
    if cfg.scope:
        data["scope"] = cfg.scope

    resp = requests.post(cfg.token_url, data=data, timeout=30)
    resp.raise_for_status()
    payload = resp.json()
    token = payload.get("access_token")
    if not token:
        raise RuntimeError(f"No access_token in response: {payload}")
    return str(token)


# ----------------------------
# Notebook-port helpers (your functions + required glue)
# ----------------------------
TOKEN_RE_DOUBLE = re.compile(r"\{\{([A-Z0-9_\/]+)\}\}")
TOKEN_RE_SINGLE = re.compile(r"\{([A-Z0-9_\/]+)\}")


def safe_num(x: Any) -> float:
    try:
        if x is None:
            return 0.0
        return float(x)
    except Exception:
        return 0.0


def remove_shape(shape) -> None:
    sp = shape._element
    sp.getparent().remove(sp)


def alpha3_to_flag(alpha3: str) -> str:
    # Optional; keep empty to avoid missing mapping issues.
    return ""


def api_get_json(base: str, token: str, path: str, params: Optional[dict] = None) -> Any:
    url = base.rstrip("/") + path
    r = requests.get(
        url,
        headers={
            "Authorization": f"Bearer {token}",
            "Accept": "application/json",
            "User-Agent": "Mozilla/5.0",
        },
        params=params or {},
        timeout=30,
    )
    if r.status_code >= 400:
        raise RuntimeError(f"GET {url} -> {r.status_code}\n{r.text[:2000]}")
    return r.json()


def items_of(obj: Any) -> List[dict]:
    if isinstance(obj, dict) and "items" in obj:
        return obj.get("items") or []
    if isinstance(obj, list):
        return obj
    return []


def pick_latest_season_ids(seasons_obj: Any, n: int = 3) -> List[int]:
    items = items_of(seasons_obj)

    def sort_key(s: dict):
        return (
            s.get("endDate") or "",
            s.get("startDate") or "",
            s.get("name") or "",
            s.get("id") or 0,
        )

    items_sorted = sorted(items, key=sort_key, reverse=True)
    out: List[int] = []
    for s in items_sorted:
        if "id" in s:
            out.append(int(s["id"]))
        if len(out) >= n:
            break
    return out


def fmt_date(iso_dt: Optional[str]) -> str:
    if not iso_dt:
        return ""
    return str(iso_dt).split("T")[0]


def prettify_camel(s: str) -> str:
    return re.sub(r"(?<!^)([A-Z])", r" \1", str(s or "")).strip()


def download_bytes(url: str, token: Optional[str] = None) -> Optional[bytes]:
    if not url:
        return None
    headers = {"User-Agent": "Mozilla/5.0"}
    if token:
        headers["Authorization"] = f"Bearer {token}"
    r = requests.get(url, headers=headers, timeout=30)
    if r.status_code >= 400:
        return None
    return r.content


def replace_textbox_exact_with_image(slide, exact_text: str, image_bytes: bytes) -> int:
    if not image_bytes:
        return 0
    replaced = 0
    for shape in list(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        t = (shape.text_frame.text or "").strip()
        if t == exact_text:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            remove_shape(shape)
            slide.shapes.add_picture(BytesIO(image_bytes), left, top, width=width, height=height)
            replaced += 1
    return replaced


def replace_tokens_in_shape(shape, values: Dict[str, str]) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False

    changed = False

    for paragraph in shape.text_frame.paragraphs:
        # 1) Try run-level replacements (preserves formatting best)
        for run in paragraph.runs:
            t = run.text or ""
            if not t:
                continue
            new_t = TOKEN_RE_DOUBLE.sub(lambda m: values.get(m.group(1), m.group(0)), t)
            new_t = TOKEN_RE_SINGLE.sub(lambda m: values.get(m.group(1), m.group(0)), new_t)
            if new_t != t:
                run.text = new_t
                changed = True

        # 2) If token spans runs, rebuild text WITHOUT clearing paragraph formatting
        combined = "".join((r.text or "") for r in paragraph.runs)
        if ("{{" in combined or "}" in combined) and paragraph.runs:
            new_combined = TOKEN_RE_DOUBLE.sub(lambda m: values.get(m.group(1), m.group(0)), combined)
            new_combined = TOKEN_RE_SINGLE.sub(lambda m: values.get(m.group(1), m.group(0)), new_combined)

            if new_combined != combined:
                # Keep paragraph properties; reuse first run, blank the rest
                paragraph.runs[0].text = new_combined
                for r in paragraph.runs[1:]:
                    r.text = ""
                changed = True

    return changed

def build_season_team_best_from_items(items: List[dict]) -> Dict[str, Dict[str, Any]]:
    """
    Returns {season_label: {"team": team_name, "mins": minutes}} picking row with max minutes.
    Skips Dutch cup: 'TOTO KNVB beker (NLD)'.
    """
    out: Dict[str, Dict[str, Any]] = {}
    for it in items or []:
        if _is_toto_knvb_beker_nld(it):
            continue

        sname = ((it.get("season") or {}).get("name") or "").strip()
        tname = ((it.get("team") or {}).get("name") or "").strip()
        if not sname or not tname:
            continue

        mins = safe_num(((it.get("stats") or {}).get("minutesPlayed")))
        mins = max(mins, safe_num(((it.get("metrics") or {}).get("minutesPlayed"))))

        prev = out.get(sname)
        if prev is None or mins > safe_num(prev.get("mins")):
            out[sname] = {"team": tname, "mins": mins}
    return out

def build_personal_values(api_base: str, token: str, player_id: int) -> Dict[str, str]:
    player = api_get_json(api_base, token, f"/api/v2/players/{player_id}")
    info = player.get("info") or {}
    team = player.get("team") or {}
    contract = player.get("contract") or {}

    seasons_obj = api_get_json(api_base, token, "/api/v2/Seasons", params={"PlayerIds": player_id, "Limit": 200})
    season_ids = pick_latest_season_ids(seasons_obj, n=5)

    # Build best team per season using fallbacks, then x-metrics override.
    cs_obj = api_get_json(
        api_base,
        token,
        "/api/v2/metrics/career-stats/players",
        params={"PlayerIds": player_id, "SeasonIds": season_ids, "Limit": 500},
    )
    season_team_best = build_season_team_best_from_items(items_of(cs_obj))

    cr_obj = api_get_json(
        api_base,
        token,
        "/api/v2/metrics/players/contribution-ratings",
        params={"PlayerIds": player_id, "SeasonIds": season_ids, "Limit": 500},
    )
    for s, v in build_season_team_best_from_items(items_of(cr_obj)).items():
        season_team_best.setdefault(s, v)

    xmetrics_obj = api_get_json(
        api_base,
        token,
        "/api/v2/metrics/players/x-metrics",
        params={"PlayerIds": player_id, "SeasonIds": season_ids, "P90": True, "Limit": 500},
    )
    for s, v in build_season_team_best_from_items(items_of(xmetrics_obj)).items():
        season_team_best[s] = v  # override if present

    # Normalize keys to "YYYY/YYYY" where possible so your CLUB_2024/2025 lookup works.
    normalized_season_team_best: Dict[str, Dict[str, Any]] = {}
    for raw_sname, payload in season_team_best.items():
        norm = normalize_season_label(raw_sname) or str(raw_sname).strip()
        if norm:
            normalized_season_team_best[norm] = payload
    season_team_best = normalized_season_team_best

    sciskill_obj = api_get_json(
        api_base,
        token,
        "/api/v2/metrics/players/sciskill",
        params={"PlayerIds": player_id, "Limit": 50},
    )
    sc = (items_of(sciskill_obj)[:1] or [{}])[0]

    roles_obj = api_get_json(
        api_base,
        token,
        "/api/v2/metrics/players/roles",
        params={"PlayerIds": player_id, "SeasonIds": season_ids, "Limit": 200},
    )
    roles_items = items_of(roles_obj)
    primary_role_row = max(roles_items, key=lambda r: r.get("minutesPlayed") or 0) if roles_items else {}
    roles = primary_role_row.get("roles") or []
    roles_sorted = sorted(roles, key=lambda rr: rr.get("fit") or 0, reverse=True)
    top_roles = ", ".join([rr.get("role") for rr in roles_sorted[:3] if rr.get("role")])

    nats = info.get("nationalities") or []
    nat_name = (nats[0].get("name") if nats else None) or (info.get("birthCountry") or {}).get("name") or ""
    nat_alpha3 = (nats[0].get("alpha3Code") if nats else None) or (info.get("birthCountry") or {}).get("alpha3Code") or ""
    flag = alpha3_to_flag(nat_alpha3)
    abb_nat = f"{flag}⎟{nat_alpha3}" if flag and nat_alpha3 else (nat_alpha3 or "")

    pos_list = info.get("positions") or []
    main_position = prettify_camel(pos_list[0]) if len(pos_list) >= 1 else ""
    sec_position = prettify_camel(pos_list[1]) if len(pos_list) >= 2 else ""

    height_m = ""
    if info.get("height") is not None:
        try:
            height_m = f"{float(info['height'])/100:.2f}".rstrip("0").rstrip(".")
        except Exception:
            height_m = ""

    contract_end = fmt_date(contract.get("contractEnd")) if contract.get("contractEnd") else ""
    on_loan_until = fmt_date(contract.get("onLoanUntil")) if contract.get("onLoanUntil") else ""
    contract_text = contract_end or on_loan_until or ("Free agent" if contract.get("isFreeAgent") else "")

    values: Dict[str, str] = {
        "PLAYER_NAME": info.get("footballName") or info.get("name") or "",
        "NATIONALITY": nat_name,
        "BIRTH_DATE": fmt_date(info.get("birthDate")),
        "HEIGHT_M": height_m,
        "CONTRACT": contract_text,
        "IS_EU": "Yes" if bool(info.get("isEuCitizen")) else "No",
        "AGENT": "",
        "PREFERRED_FOOT": info.get("preferredFoot") or "",
        "MAIN_POSITION": main_position,
        "SEC_POSITION": sec_position,
        "SCISKILL": str(sc.get("sciskill", "")).strip(),
        "POTENTIAL": str(sc.get("potential", "")).strip(),
        "SCISKILL_DEV_6M": str(sc.get("sciskillDevelopmentSixMonths", "")).strip(),
        "POTENTIAL_DEV_6M": str(sc.get("potentialDevelopmentSixMonths", "")).strip(),
        "TOP_ROLES": top_roles,
        "PRIMARY_ROLE_POSITION": prettify_camel(primary_role_row.get("position") or ""),
        "ABB_NATIONALITY": abb_nat,

        # Clubs (now includes 25/26 + 22/23)
        "CLUB_2025/2026": (season_team_best.get("2025/2026") or {}).get("team", ""),
        "CLUB_2024/2025": (season_team_best.get("2024/2025") or {}).get("team", ""),
        "CLUB_2023/2024": (season_team_best.get("2023/2024") or {}).get("team", ""),
        "CLUB_2022/2023": (season_team_best.get("2022/2023") or {}).get("team", ""),
      
        "IMAGE": "{IMAGE}",
        "PRESTATIES_FIGURE": "{PRESTATIES_FIGURE}",
        "_POSITIONS_ORDERED": pos_list,
    }

  # Fallback: if 2022/2023 club is empty, copy 2023/2024 club
    if not (values.get("CLUB_2022/2023") or "").strip():
      values["CLUB_2022/2023"] = (values.get("CLUB_2023/2024") or "").strip()


    values["_PLAYER_IMAGE_URL"] = info.get("imageUrl") or ""
    values["_TEAM_IMAGE_URL"] = team.get("imageUrl") or ""
    return values

def apply_season_row_tokens_blank_if_missing(
    values: Dict[str, str],
    season_label: str,
    season_id_by_label: Dict[str, int],
    stats_by_sid: Dict[int, Dict[str, int]],
    club_key: str,
    g_key: str,
    m_key: str,
    go_key: str,
    a_key: str,
) -> None:
    sid = season_id_by_label.get(season_label)
    stats = stats_by_sid.get(sid) if isinstance(sid, int) else None
    if not sid or stats is None:
        values[club_key] = ""
        values[g_key] = ""
        values[m_key] = ""
        values[go_key] = ""
        values[a_key] = ""
        return

    values[g_key] = str(stats["GAMES"])
    values[m_key] = str(stats["MINUTES"])
    values[go_key] = str(stats["GOALS"])
    values[a_key] = str(stats["ASSISTS"])

def _normalize_season_name(name: str) -> str:
    return normalize_season_label(name) if "normalize_season_label" in globals() else str(name).strip()


def build_season_ids_by_label(seasons_obj: dict) -> Dict[str, List[int]]:
    """
    Returns mapping like {"2024/2025": [123, 456, ...], ...}
    because SciSports has multiple season IDs per year (per competition).
    """
    out: Dict[str, List[int]] = {}
    for it in items_of(seasons_obj):
        sid = it.get("id")
        sname = _normalize_season_name((it.get("name") or "").strip())
        if isinstance(sid, int) and sname:
            out.setdefault(sname, []).append(sid)
    return out


def _find_stat(stats: dict, keys: List[str]) -> float:
    for k in keys:
        if k in stats:
            return safe_num(stats[k])
    lower = {str(k).lower(): k for k in stats.keys()}
    for k in keys:
        lk = k.lower()
        if lk in lower:
            return safe_num(stats[lower[lk]])
    return 0.0

def get_teamwise_stats_by_season(
    api_base: str,
    token: str,
    player_id: int,
    season_ids: List[int],
) -> Dict[int, Dict[str, Dict[str, int]]]:
    """
    Returns:
      { season_id: { team_name: {"GAMES":..,"MINUTES":..,"GOALS":..,"ASSISTS":..} } }

    Rule:
    - Sums across ALL competitions within a season.
    - Aggregates separately per team.
    - No KNVB-beker filtering here (you asked to include all competitions).
    """
    # Games + minutes from contribution-ratings
    cr = api_get_json(
        api_base,
        token,
        "/api/v2/metrics/players/contribution-ratings",
        params={"PlayerIds": player_id, "SeasonIds": season_ids, "Limit": 500},
    )
    cr_items = cr.get("items") or []

    out: Dict[int, Dict[str, Dict[str, int]]] = {}

    def _ensure(sid: int, team: str) -> Dict[str, int]:
        out.setdefault(sid, {})
        out[sid].setdefault(team, {"GAMES": 0, "MINUTES": 0, "GOALS": 0, "ASSISTS": 0})
        return out[sid][team]

    for it in cr_items:
        sid = (it.get("season") or {}).get("id")
        if not isinstance(sid, int) or sid not in season_ids:
            continue

        team_name = ((it.get("team") or {}).get("name") or "").strip() or "Unknown"
        stats = it.get("stats") or {}

        bucket = _ensure(sid, team_name)
        bucket["GAMES"] += int(round(safe_num(stats.get("matchesPlayed"))))
        bucket["MINUTES"] += int(round(safe_num(stats.get("minutesPlayed"))))

    # Goals + assists from career-stats
    cs = api_get_json(
        api_base,
        token,
        "/api/v2/metrics/career-stats/players",
        params={"PlayerIds": player_id, "SeasonIds": season_ids, "Limit": 500},
    )
    cs_items = cs.get("items") or []

    for it in cs_items:
        sid = (it.get("season") or {}).get("id") or it.get("seasonId")
        if not isinstance(sid, int) or sid not in season_ids:
            continue

        team_name = ((it.get("team") or {}).get("name") or "").strip() or "Unknown"
        stats = it.get("stats") or {}

        goals = int(round(_find_stat(stats, ["goal", "goals", "goalNonPenalty", "goal_non_penalty"])))
        assists = int(round(_find_stat(stats, ["assist", "assists"])))

        bucket = _ensure(sid, team_name)
        bucket["GOALS"] += goals
        bucket["ASSISTS"] += assists

    return out

def get_career_stats_totals_by_season_team(
    api_base: str,
    token: str,
    player_id: int,
    season_ids: List[int],
) -> Dict[int, Dict[str, Dict[str, int]]]:
    """
    Returns: {season_id: {team_name: {"GAMES","MINUTES","GOALS","ASSISTS"}}}

    Fixes double counting by:
    1) If a "Total" row exists for a season+team -> use that only.
    2) Otherwise, for each season+team+competition take MAX (rows are often cumulative),
       then sum across competitions.
    """

    cs = api_get_json(
        api_base,
        token,
        "/api/v2/metrics/career-stats/players",
        params={"PlayerIds": player_id, "SeasonIds": season_ids, "Limit": 2000},
    )
    items = items_of(cs)

    # helper: read comp identity (id if possible, else name, else unknown)
    def comp_key(it: dict) -> str:
        comp = it.get("competition") or it.get("competitionGroup") or {}
        cid = comp.get("id")
        cname = (comp.get("name") or "").strip()
        if cid is not None:
            return f"id:{cid}"
        if cname:
            return f"name:{cname.lower()}"
        # sometimes nested deeper
        cname2 = ((it.get("competition") or {}).get("name") or "").strip()
        return f"name:{cname2.lower()}" if cname2 else "unknown"

    def is_total_row(it: dict) -> bool:
        comp = it.get("competition") or it.get("competitionGroup") or {}
        cname = (comp.get("name") or "").strip().lower()
        # common patterns; adjust if your payload uses other labels
        return cname in {"total", "all", "overall", "all competitions", "all competitions total"}

    # Step 1: collect best "Total row" per (sid, team) if present
    best_total: Dict[Tuple[int, str], Dict[str, int]] = {}

    # Step 2: otherwise collect MAX per (sid, team, competition)
    max_per_comp: Dict[Tuple[int, str, str], Dict[str, int]] = {}

    for it in items:
        sid = (it.get("season") or {}).get("id") or it.get("seasonId")
        # inside loop, before summing:
        comp = (it.get("competition") or {}).get("name") or (it.get("league") or {}).get("name") or ""
        season_name = (it.get("season") or {}).get("name") or ""
        # print(comp, season_name, team_name, games, minutes, goals, assists)
        if not isinstance(sid, int) or sid not in season_ids:
            continue

        team_name = ((it.get("team") or {}).get("name") or "").strip() or "Unknown"
        stats = it.get("stats") or {}
        metrics = it.get("metrics") or {}

        games = int(round(_find_stat(stats, ["matchesPlayed", "matchPlayed", "matches", "games"])))
        minutes = int(round(
            max(
                _find_stat(stats, ["minutesPlayed", "minutes"]),
                _find_stat(metrics, ["minutesPlayed", "minutes"]),
            )
        ))
        goals = int(round(_find_stat(stats, ["goal", "goals", "goalNonPenalty", "goal_non_penalty"])))
        assists = int(round(_find_stat(stats, ["assist", "assists"])))

        row = {"GAMES": games, "MINUTES": minutes, "GOALS": goals, "ASSISTS": assists}

        if is_total_row(it):
            key2 = (sid, team_name)
            prev = best_total.get(key2)
            # keep the best total row (usually highest minutes)
            if prev is None or row["MINUTES"] > prev["MINUTES"]:
                best_total[key2] = row
            continue

        ck = comp_key(it)
        key3 = (sid, team_name, ck)
        prev = max_per_comp.get(key3)
        if prev is None:
            max_per_comp[key3] = row
        else:
            # rows are often cumulative -> take MAX, not SUM
            prev["GAMES"] = max(prev["GAMES"], row["GAMES"])
            prev["MINUTES"] = max(prev["MINUTES"], row["MINUTES"])
            prev["GOALS"] = max(prev["GOALS"], row["GOALS"])
            prev["ASSISTS"] = max(prev["ASSISTS"], row["ASSISTS"])

    # Build final output
    out: Dict[int, Dict[str, Dict[str, int]]] = {}

    # If total rows exist, they win
    for (sid, team), row in best_total.items():
        out.setdefault(sid, {})
        out[sid][team] = dict(row)

    # Otherwise sum MAX-per-competition
    for (sid, team, _ck), row in max_per_comp.items():
        # don't overwrite total rows
        if (sid, team) in best_total:
            continue
        out.setdefault(sid, {})
        out[sid].setdefault(team, {"GAMES": 0, "MINUTES": 0, "GOALS": 0, "ASSISTS": 0})
        out[sid][team]["GAMES"] += row["GAMES"]
        out[sid][team]["MINUTES"] += row["MINUTES"]
        out[sid][team]["GOALS"] += row["GOALS"]
        out[sid][team]["ASSISTS"] += row["ASSISTS"]

    return out



def get_games_minutes_goals_assists_by_season(
    api_base: str,
    token: str,
    player_id: int,
    season_ids: List[int],
) -> Dict[int, Dict[str, int]]:
    cr = api_get_json(
        api_base,
        token,
        "/api/v2/metrics/players/contribution-ratings",
        params={"PlayerIds": player_id, "SeasonIds": season_ids, "Limit": 500},
    )
    cr_items = cr.get("items") or []

    gm_by_sid: Dict[int, Dict[str, int]] = {}
    for it in cr_items:
        if _is_toto_knvb_beker_nld(it):
            continue
        sid = (it.get("season") or {}).get("id")
        if not isinstance(sid, int):
            continue
        stats = it.get("stats") or {}
        gm_by_sid.setdefault(sid, {"GAMES": 0, "MINUTES": 0})
        gm_by_sid[sid]["GAMES"] += int(round(safe_num(stats.get("matchesPlayed"))))
        gm_by_sid[sid]["MINUTES"] += int(round(safe_num(stats.get("minutesPlayed"))))

    cs = api_get_json(
        api_base,
        token,
        "/api/v2/metrics/career-stats/players",
        params={"PlayerIds": player_id, "SeasonIds": season_ids, "Limit": 500},
    )
    cs_items = cs.get("items") or []

    ga_by_sid: Dict[int, Dict[str, int]] = {}
    for it in cs_items:
        if _is_toto_knvb_beker_nld(it):
            continue
        sid = (it.get("season") or {}).get("id") or it.get("seasonId")
        if not isinstance(sid, int):
            continue
        stats = it.get("stats") or {}
        goals = int(round(_find_stat(stats, ["goal", "goals", "goalNonPenalty", "goal_non_penalty"])))
        assists = int(round(_find_stat(stats, ["assist", "assists"])))

        prev = ga_by_sid.get(sid)
        if prev is None:
            ga_by_sid[sid] = {"GOALS": goals, "ASSISTS": assists}
        else:
            ga_by_sid[sid]["GOALS"] = max(prev["GOALS"], goals)
            ga_by_sid[sid]["ASSISTS"] = max(prev["ASSISTS"], assists)

    out: Dict[int, Dict[str, int]] = {}
    for sid in season_ids:
        if sid not in gm_by_sid and sid not in ga_by_sid:
            continue
        out[sid] = {
            "GAMES": int(gm_by_sid.get(sid, {}).get("GAMES", 0)),
            "MINUTES": int(gm_by_sid.get(sid, {}).get("MINUTES", 0)),
            "GOALS": int(ga_by_sid.get(sid, {}).get("GOALS", 0)),
            "ASSISTS": int(ga_by_sid.get(sid, {}).get("ASSISTS", 0)),
        }
    return out
def apply_season_row_tokens_teamwise(
    values: Dict[str, str],
    season_label: str,
    season_ids_by_label: Dict[str, List[int]],
    totals_by_sid: Dict[int, Dict[str, Dict[str, int]]],
    club_key: str,
    g_key: str,
    m_key: str,
    go_key: str,
    a_key: str,
) -> None:
    season_ids = season_ids_by_label.get(season_label) or []
    if not season_ids:
        values[club_key] = values[g_key] = values[m_key] = values[go_key] = values[a_key] = ""
        return

    # Merge across ALL competition-season IDs for that label
    merged: Dict[str, Dict[str, int]] = {}
    for sid in season_ids:
        teams = totals_by_sid.get(sid) or {}
        for team_name, st in teams.items():
            bucket = merged.setdefault(team_name, {"GAMES": 0, "MINUTES": 0, "GOALS": 0, "ASSISTS": 0})
            bucket["GAMES"] += int(st.get("GAMES", 0))
            bucket["MINUTES"] += int(st.get("MINUTES", 0))
            bucket["GOALS"] += int(st.get("GOALS", 0))
            bucket["ASSISTS"] += int(st.get("ASSISTS", 0))

    if not merged:
        values[club_key] = values[g_key] = values[m_key] = values[go_key] = values[a_key] = ""
        return

    ordered = sorted(merged.items(), key=lambda kv: kv[1].get("MINUTES", 0), reverse=True)
    sep = " / " if len(ordered) > 1 else ""

    values[club_key] = sep.join([team for team, _ in ordered])
    values[g_key] = sep.join([str(st["GAMES"]) for _, st in ordered])
    values[m_key] = sep.join([str(st["MINUTES"]) for _, st in ordered])
    values[go_key] = sep.join([str(st["GOALS"]) for _, st in ordered])
    values[a_key] = sep.join([str(st["ASSISTS"]) for _, st in ordered])

def apply_season_row_tokens(
    values: Dict[str, str],
    season_label: str,
    season_id_by_label: Dict[str, int],
    stats_by_sid: Dict[int, Dict[str, int]],
    club_team_by_label: Dict[str, Dict[str, Any]],
    club_key: str,
    g_key: str,
    m_key: str,
    go_key: str,
    a_key: str,
) -> None:
    """
    If season missing -> blanks CLUB + all stats placeholders.
    """
    sid = season_id_by_label.get(season_label)
    stats = stats_by_sid.get(sid) if isinstance(sid, int) else None

    if not sid or stats is None:
        values[club_key] = ""
        values[g_key] = ""
        values[m_key] = ""
        values[go_key] = ""
        values[a_key] = ""
        return

    club = (club_team_by_label.get(season_label) or {}).get("team", "")
    values[club_key] = str(club or "")
    values[g_key] = str(stats["GAMES"])
    values[m_key] = str(stats["MINUTES"])
    values[go_key] = str(stats["GOALS"])
    values[a_key] = str(stats["ASSISTS"])


def compute_strengths_and_percentile_from_api(
    api_base: str,
    token: str,
    player_id: int,
    season_id: int,
) -> Tuple[str, float]:
    xm = api_get_json(
        api_base,
        token,
        "/api/v2/metrics/players/x-metrics",
        params={"PlayerIds": player_id, "SeasonIds": [season_id], "Limit": 200},
    )
    items = xm.get("items") or []
    best = None
    best_min = -1
    for it in items:
        mins = safe_num(((it.get("metrics") or {}).get("minutesPlayed")))
        if mins > best_min:
            best_min = mins
            best = it

    if not best:
        return ("High performance on Crossing, Passing, Offensive positioning", 50.0)

    m = best.get("metrics") or {}
    crossing = safe_num(m.get("xaCross"))
    passing = max(safe_num(m.get("xaOpenPlay")), safe_num(m.get("xa")))
    off_pos = max(safe_num(m.get("xgOpenPlay")), safe_num(m.get("xgShotsClose")))

    triples = [("Crossing", crossing), ("Passing", passing), ("Offensive positioning", off_pos)]
    triples_sorted = sorted(triples, key=lambda t: t[1], reverse=True)
    strengths = ", ".join([t[0] for t in triples_sorted[:3]])
    strengths_line = f"High performance on {strengths}"

    xg_plus_xa = safe_num(m.get("xgPlusXa"))
    minutes = max(1.0, safe_num(m.get("minutesPlayed")))
    per90 = xg_plus_xa / (minutes / 90.0)
    percentile = 100.0 / (1.0 + math.exp(-2.0 * (per90 - 0.30)))
    percentile = max(0.0, min(100.0, percentile))
    return strengths_line, float(percentile)


def generate_performance_plot_simple(
    out_png: str,
    title: str,
    subtitle: str,
    player_name: str,
    strengths_line: str,
    percentile: float,
) -> None:
    fig = plt.figure(figsize=(8.0, 5.0), dpi=150)
    ax = fig.add_axes([0.08, 0.22, 0.84, 0.52])
    ax.set_xlim(0, 100)
    ax.set_ylim(-1, 1)
    ax.fill_between([0, 33], -0.25, 0.25, alpha=0.15)
    ax.fill_between([33, 66], -0.25, 0.25, alpha=0.15)
    ax.fill_between([66, 100], -0.25, 0.25, alpha=0.15)
    ax.scatter([percentile], [0], s=250)
    ax.set_yticks([])
    ax.set_xticks([0, 33, 66, 100])
    ax.set_xlabel("Percentile vs peers")
    for spine in ax.spines.values():
        spine.set_visible(False)
    fig.text(0.08, 0.92, title, fontsize=18, weight="bold")
    fig.text(0.08, 0.875, subtitle, fontsize=12)
    fig.text(0.08, 0.80, player_name, fontsize=14, weight="bold")
    fig.text(0.08, 0.08, strengths_line, fontsize=12)
    fig.savefig(out_png, bbox_inches="tight")
    plt.close(fig)


# ----------------------------
# Radar helpers (unchanged)
# ----------------------------
def _split_label_unit(label: str) -> Tuple[str, str]:
    if label.endswith("(m)"):
        return label.replace(" (m)", ""), "m"
    return label, ""


def _norm_name(s: str) -> str:
    if s is None:
        return ""
    s = str(s).strip().lower()
    s = unicodedata.normalize("NFKD", s)
    s = "".join(ch for ch in s if not unicodedata.combining(ch))
    s = re.sub(r"[^a-z0-9\s]", " ", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s


def _similarity(a: str, b: str) -> float:
    return SequenceMatcher(None, a, b).ratio()


def pick_player_row_by_name(df_bench: pd.DataFrame, player_name: str, team_name: str | None = None) -> pd.Series:
    if "player" not in df_bench.columns:
        raise ValueError("bench.csv must contain a 'player' column.")

    target = _norm_name(player_name)
    if not target:
        raise ValueError("Empty player_name passed to matcher.")

    df = df_bench.copy()
    df["_player_norm"] = df["player"].astype(str).map(_norm_name)

    exact = df[df["_player_norm"] == target]
    if len(exact) == 0:
        exact = df[df["_player_norm"].str.contains(target, na=False)]

    if team_name and "team" in df.columns and len(exact) > 1:
        tnorm = _norm_name(team_name)
        exact_team = exact[exact["team"].astype(str).map(_norm_name) == tnorm]
        if len(exact_team) > 0:
            exact = exact_team

    if len(exact) > 0:
        if "total_minutes" in exact.columns:
            exact = exact.sort_values("total_minutes", ascending=False)
        return exact.iloc[0]

    sims = df["_player_norm"].map(lambda n: _similarity(target, n))
    best_idx = sims.idxmax()
    best_score = float(sims.loc[best_idx])
    if best_score < 0.80:
        top = df.assign(_sim=sims).sort_values("_sim", ascending=False).head(5)[["player", "_sim"]]
        raise ValueError(
            f"No reliable name match for '{player_name}'. Best similarity={best_score:.3f}. "
            f"Top candidates:\n{top.to_string(index=False)}"
        )

    best_name = df.loc[best_idx, "_player_norm"]
    near = df[df["_player_norm"] == best_name]
    if team_name and "team" in near.columns and len(near) > 1:
        tnorm = _norm_name(team_name)
        near_team = near[near["team"].astype(str).map(_norm_name) == tnorm]
        if len(near_team) > 0:
            near = near_team

    if "total_minutes" in near.columns:
        near = near.sort_values("total_minutes", ascending=False)
    return near.iloc[0]


def load_bench_csv(path: str) -> pd.DataFrame:
    if not os.path.exists(path):
        raise FileNotFoundError(f"Missing {path}. Put bench.csv in the repo root.")
    df = pd.read_csv(path)
    if "player" not in df.columns:
        raise ValueError("bench.csv must include a 'player' column.")
    return df


def compute_maxes(
    df_bench: pd.DataFrame,
    player_row: pd.Series,
    labels: List[str],
    use_custom: Optional[Dict[str, float]],
) -> Dict[str, float]:
    if use_custom:
        return {lab: float(use_custom[lab]) for lab in labels}

    subset = df_bench
    if "position_group" in df_bench.columns and pd.notna(player_row.get("position_group")):
        subset = df_bench[df_bench["position_group"] == player_row["position_group"]]
        if len(subset) < 30:
            subset = df_bench

    maxes: Dict[str, float] = {}
    for lab in labels:
        col = RADAR_METRICS_MAP[lab]
        s = pd.to_numeric(subset[col], errors="coerce").dropna()
        if len(s) == 0:
            maxes[lab] = 1.0
        else:
            mx = float(s.quantile(0.95))
            maxes[lab] = mx if mx > 0 else 1.0
    return maxes
  
_compute_maxes = compute_maxes


def generate_radar_chart_for_player(
    df_bench: pd.DataFrame,
    player_name: str,
    out_png: str,
    custom_maxes: Optional[Dict[str, float]] = None,
    team_name: Optional[str] = None,
) -> Dict[str, float]:
    """
    Generates the radar chart PNG with 50% background transparency.

    Add-ons:
    - Orange dashed line (no fill): KKD average within same position_group (fallback: all KKD).
    - Legend.
    - Maxes: per metric = highest KKD player value, rounded up to nearest 100 (unless custom_maxes provided).
    """

    row = pick_player_row_by_name(df_bench, player_name=player_name, team_name=team_name)

    labels = list(RADAR_METRICS_MAP.keys())

    # -------------------------
    # Player raw values
    # -------------------------
    raw_vals = []
    for lab in labels:
        col = RADAR_METRICS_MAP[lab]
        v = float(pd.to_numeric(row.get(col, np.nan), errors="coerce"))
        if math.isnan(v):
            v = 0.0
        raw_vals.append(v)
    raw_vals = np.array(raw_vals, dtype=float)

    # -------------------------
    # Maxes (KKD max rounded to 100)
    # -------------------------
    # Maxes = highest value in df_bench per metric (rounded to whole number)
    # -------------------------
    if custom_maxes:
        maxes_dict = {lab: float(custom_maxes[lab]) for lab in labels}
    else:
        maxes_dict = {}
        for lab in labels:
            col = RADAR_METRICS_MAP[lab]
            s = pd.to_numeric(df_bench[col], errors="coerce").dropna() if col in df_bench.columns else pd.Series([], dtype=float)
            mx = float(s.max()) if len(s) else 0.0
            maxes_dict[lab] = max(1.0, float(int(round(mx))))

    max_vals = np.array([maxes_dict[lab] for lab in labels], dtype=float)
    max_vals = np.where(max_vals <= 0, 100.0, max_vals)

    # -------------------------
    # Normalized player values
    # -------------------------
    norm = np.clip(raw_vals / max_vals, 0.0, 1.0)
    norm_closed = np.r_[norm, norm[0]]

    # -------------------------
    # KKD average (same position_group; fallback: all KKD)
    # -------------------------
    kkd_norm_closed = None
    if "division" in df_bench.columns:
        kkd_df = df_bench[df_bench["division"].astype(str) == "KKD"]

        if "position_group" in df_bench.columns and pd.notna(row.get("position_group")):
            kkd_pg = kkd_df[kkd_df["position_group"] == row["position_group"]]
            if len(kkd_pg) > 0:
                kkd_df = kkd_pg

        if len(kkd_df) > 0:
            kkd_avg_vals = []
            for lab in labels:
                col = RADAR_METRICS_MAP[lab]
                s = pd.to_numeric(kkd_df[col], errors="coerce").dropna() if col in kkd_df.columns else pd.Series([], dtype=float)
                kkd_avg_vals.append(float(s.mean()) if len(s) else 0.0)

            kkd_avg_vals = np.array(kkd_avg_vals, dtype=float)
            kkd_norm = np.clip(kkd_avg_vals / max_vals, 0.0, 1.0)
            kkd_norm_closed = np.r_[kkd_norm, kkd_norm[0]]

    # -------------------------
    # Plot geometry
    # -------------------------
    N = len(labels)
    angles = np.linspace(0, 2 * np.pi, N, endpoint=False)
    angles_closed = np.r_[angles, angles[0]]

    fig = plt.figure(figsize=(7.8, 7.8), dpi=150)
    ax = plt.subplot(111, polar=True)
    ax.set_aspect("equal", adjustable="box")

    ax.set_theta_offset(np.pi / 2)
    ax.set_theta_direction(-1)

    # Player polygon
    player_line, = ax.plot(angles_closed, norm_closed, linewidth=3, marker="o", markersize=6)
    ax.fill(angles_closed, norm_closed, alpha=0.20)

    # KKD dashed overlay (no fill)
    kkd_line = None
    if kkd_norm_closed is not None:
        kkd_line, = ax.plot(
            angles_closed,
            kkd_norm_closed,
            linewidth=2.5,
            linestyle="--",
            color="orange",
        )

    # Legend
    handles = [player_line]
    legend_labels = [player_name]
    if kkd_line is not None:
        handles.append(kkd_line)
        legend_labels.append("KKD avg (same position)")
    ax.legend(
        handles,
        legend_labels,
        loc="upper right",
        bbox_to_anchor=(1.12, 1.08),
        frameon=False,
        fontsize=10,
    )

    # -------------------------
    # Rings / labels
    # -------------------------
    rings = [0.25, 0.50, 0.75, 1.00]
    ax.set_ylim(0, 1.0)
    ax.set_yticks(rings)
    ax.set_yticklabels([""] * len(rings))
    ax.yaxis.grid(True, linewidth=1)
    ax.xaxis.grid(True, linewidth=1)

    ax.set_xticks(angles)
    ax.set_xticklabels([""] * N)

    tick_fontsize = 8
    for ang, mx, lab in zip(angles, max_vals, labels):
        _, unit = _split_label_unit(lab)
        for r in rings:
            v = r * mx
            txt = f"{v:.0f}" if unit == "" else f"{v:.0f} {unit}"
            ax.text(ang, r, txt, fontsize=tick_fontsize, ha="center", va="center")

    R_VALUE = 1.15  # keep default for all
    TOTAL_DISTANCE_VALUE_R = 1.08
    TOTAL_DISTANCE_NAME_R = 1.15
    
    for lab, ang, val in zip(labels, angles, raw_vals):
        name, unit = _split_label_unit(lab)
        a = (ang + np.pi / 2) % (2 * np.pi)
        c = np.cos(a)
        if c > 0.25:
            ha = "left"
        elif c < -0.25:
            ha = "right"
        else:
            ha = "center"
    
        # default positions
        r_name = NAME_R_MULT.get(lab, DEFAULT_NAME_R)
        r_value = R_VALUE
    
        # only adjust Total distance
        if lab == "Total distance (m)":
            r_name = TOTAL_DISTANCE_NAME_R
            r_value = TOTAL_DISTANCE_VALUE_R
    
        ax.text(ang, r_name, name, fontsize=13, fontweight="bold", ha=ha, va="center")
        val_line = f"{val:.0f}" if unit == "" else f"{val:.0f} {unit}"
        ax.text(ang, r_value, val_line, fontsize=13, fontweight="bold", ha=ha, va="center")


    # --- FINAL PLOT STYLING (KEEP YOUR 50% TRANSPARENCY LOGIC) ---
    ax.spines["polar"].set_visible(False)

    fig.patch.set_facecolor((1, 1, 1, 0.5))   # rectangle background
    ax.patch.set_facecolor((1, 1, 1, 0.5))    # circle background

    # Symmetric margins -> keeps the radar centered in the exported image
    fig.subplots_adjust(left=0.12, right=0.88, top=0.95, bottom=0.05)
    fig.savefig(out_png, transparent=False)

    plt.close(fig)

    return {lab: float(v) for lab, v in zip(labels, raw_vals)}


def get_local_player_image_path(player_name: str, photos_dir: str) -> Optional[str]:
    if not player_name or not os.path.isdir(photos_dir):
        return None

    candidates = [
        f"{player_name}.png",
        f"{player_name}.jpg",
        f"{player_name}.jpeg",
        f"{player_name.strip()}.png",
    ]

    for name in candidates:
        p = os.path.join(photos_dir, name)
        if os.path.exists(p):
            return p

    # case-insensitive fallback
    target_low = f"{player_name.lower().strip()}.png"
    for f in os.listdir(photos_dir):
        if f.lower() == target_low:
            return os.path.join(photos_dir, f)

    return None



# ----------------------------
# PPTX image insertion
# ----------------------------
def insert_image_at_token_exact(slide, token: str, image_path: str) -> int:
    """
    Replace a shape whose entire text equals token with an image at same position.
    """
    replaced = 0
    for shape in list(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        if (shape.text_frame.text or "").strip() != token:
            continue
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        remove_shape(shape)
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)
        replaced += 1
    return replaced
from pptx.dml.color import RGBColor

POSITION_TO_NUMBER: Dict[str, int] = {
    # GK
    "Goalkeeper": 1,

    # Back line
    "RightBack": 2,
    "RightFullback": 2,
    "Right Back": 2,
    "Centre Back": 3,
    "CentreBack": 3,
    "LeftBack": 4,
    "Left Back": 4,

    # Midfield (common API variants)
    "DefensiveMidfield": 6,
    "Defensive Midfield": 6,
    "CentreMidfield": 8,
    "Centre Midfield": 8,
    "AttackingMidfield": 10,
    "Attacking Midfield": 10,

    # Wings / forwards
    "RightWing": 7,
    "Right Wing": 7,
    "Left Wing": 11,
    "LeftWing": 11,
    "Striker": 9,
    "CentreForward": 9,
    "Centre Forward": 9,
}


MAIN_BLUE = RGBColor(0, 83, 159)      # adjust to your exact template blue if needed
SECOND_BLUE = RGBColor(0, 142, 204)   # adjust to your exact template light-blue if needed
def apply_position_coloring(slide, ordered_positions: List[str]) -> None:
    if not ordered_positions:
        return

    main_num = POSITION_TO_NUMBER.get(ordered_positions[0])
    secondary_nums = [
        POSITION_TO_NUMBER.get(p)
        for p in ordered_positions[1:3]
        if p in POSITION_TO_NUMBER
    ]

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        txt = (shape.text_frame.text or "").strip()
        if not txt.isdigit():
            continue

        num = int(txt)
        if main_num is not None and num == main_num:
            shape.fill.solid()
            shape.fill.fore_color.rgb = MAIN_BLUE
        elif num in secondary_nums:
            shape.fill.solid()
            shape.fill.fore_color.rgb = SECOND_BLUE

# ----------------------------
# Template filling (full notebook-style flow)
# ----------------------------

def fill_template_full(
    template_path: str,
    out_pptx_path: str,
    df_bench: pd.DataFrame,
    player_name_ui: str,
    player_id: int,
    api_base: str,
    token: str,
    performance_upload_bytes: Optional[bytes],
) -> Dict[str, Any]:
    inserted = {"player_image": 0, "radar": 0, "performance": 0, "text_shapes_changed": 0}

    if not os.path.exists(template_path):
        raise FileNotFoundError(...)

    if not (api_base and token and player_id):
        raise RuntimeError("Missing api_base/token/player_id - cannot build full report values.")

    # ✅ Always define values here, unconditionally
    values: Dict[str, Any] = build_personal_values(api_base, token, int(player_id))

    # ✅ Ensure _POSITIONS_ORDERED exists AFTER values exists
    values.setdefault("_POSITIONS_ORDERED", [])

    prs = Presentation(template_path)

    # Build values (old logic)
    # Seasons (keep what you already have)
    seasons_obj = api_get_json(api_base, token, "/api/v2/Seasons", params={"PlayerIds": player_id, "Limit": 500})
    season_ids_latest5 = pick_latest_season_ids(seasons_obj, n=5)
    season_ids_by_label = build_season_ids_by_label(seasons_obj)
    
    target_labels = ["2025/2026", "2024/2025"]
    target_season_ids = []
    for lbl in target_labels:
        target_season_ids.extend(season_ids_by_label.get(lbl, []))

    
    totals_by_sid = get_career_stats_totals_by_season_team(
        api_base=api_base,
        token=token,
        player_id=player_id,
        season_ids=target_season_ids,
    )
    
    apply_season_row_tokens_teamwise(
        values=values,
        season_label="2025/2026",
        season_ids_by_label=season_ids_by_label,
        totals_by_sid=totals_by_sid,
        club_key="CLUB_2025/2026",
        g_key="G25/26",
        m_key="M25/26",
        go_key="GO25/26",
        a_key="A25/26",
    )
    
    apply_season_row_tokens_teamwise(
        values=values,
        season_label="2024/2025",
        season_ids_by_label=season_ids_by_label,
        totals_by_sid=totals_by_sid,
        club_key="CLUB_2024/2025",
        g_key="G24/25",
        m_key="M24/25",
        go_key="GO24/25",
        a_key="A24/25",
    )

    
    # For older seasons: club only (already set in build_personal_values via season_team_best)
    # Do NOT override CLUB_2023/2024 and CLUB_2022/2023 here.



    # strengths_line, percentile = compute_strengths_and_percentile_from_api(
    #     api_base=api_base,
    #     token=token,
    #     player_id=int(player_id),
    #     season_id=int(latest_season_id) if latest_season_id else int(stats.get("season_id", 0)),
    # )
    # values["STRENGTHS"] = strengths_line

    # Player image bytes (local first; fallback API imageUrl)
    player_img_bytes: Optional[bytes] = None
    local_path = (
        get_local_player_image_path(values.get("PLAYER_NAME") or "", PLAYER_PHOTOS_DIR)
        or get_local_player_image_path(player_name_ui, PLAYER_PHOTOS_DIR)
    )
    if local_path:
        with open(local_path, "rb") as f:
            player_img_bytes = f.read()
    else:
        player_img_bytes = download_bytes(values.get("_PLAYER_IMAGE_URL", ""), token=None)

    # Radar chart
    radar_png = os.path.join(os.path.dirname(out_pptx_path), "radar_chart.png")
    radar_used = generate_radar_chart_for_player(
        df_bench=df_bench,
        player_name=values.get("PLAYER_NAME") or player_name_ui,
        out_png=radar_png,
        custom_maxes=CUSTOM_MAXES,  # <- enable custom maxes
        team_name=values.get("CLUB_2024/2025") or values.get("CLUB_2023/2024") or None,
    )
    

    # Performance chart: upload wins; else auto-generate like notebook
    perf_png: Optional[str] = None
    perf_used = "missing"
    
    if performance_upload_bytes:
        perf_png = os.path.join(os.path.dirname(out_pptx_path), "performance_chart.png")
        with open(perf_png, "wb") as f:
            f.write(performance_upload_bytes)
        perf_used = "uploaded"

    
        inserted = {
            "player_image": 0,
            "radar": 0,
            "performance": 0,
            "text_shapes_changed": 0,
        }

    for slide in prs.slides:
        if player_img_bytes:
            inserted["player_image"] += replace_textbox_exact_with_image(slide, "{IMAGE}", player_img_bytes)
    
        inserted["radar"] += insert_image_at_token_exact(slide, "{{RADAR_CHART}}", radar_png)
    
        if perf_png:
            inserted["performance"] += insert_image_at_token_exact(slide, "{{PERFORMANCE_CHART}}", perf_png)
            inserted["performance"] += replace_textbox_exact_with_image(slide, "{PRESTATIES_FIGURE}", performance_upload_bytes)
    
    for shape in slide.shapes:
        if replace_tokens_in_shape(shape, values):
            inserted["text_shapes_changed"] += 1
    
    apply_position_coloring(slide, values.get("_POSITIONS_ORDERED", []))
    

    prs.save(out_pptx_path)

    return {
        "player_name_ui": player_name_ui,
        "player_id": player_id,
        "values_keys": sorted([k for k in values.keys() if not k.startswith("_")]),
        "radar_used": radar_used,
        "perf_used": perf_used,
        "inserted": inserted,
    }


# ----------------------------
# PDF conversion (best-effort)
# ----------------------------
def can_convert_to_pdf() -> bool:
    return shutil.which("soffice") is not None


def convert_pptx_to_pdf(pptx_path: str, out_dir: str) -> str:
    if shutil.which("soffice") is None:
        raise RuntimeError("LibreOffice (soffice) not found on PATH.")

    # Dedicated LO profile to avoid first-run + font cache weirdness on Streamlit Cloud
    profile_dir = Path(out_dir) / "lo_profile"
    profile_dir.mkdir(parents=True, exist_ok=True)
    profile_uri = profile_dir.resolve().as_uri()

    # Force Impress PDF export filter + disable image downsampling + lossless compression
    convert_to = (
        'pdf:impress_pdf_Export:'
        '{"ReduceImageResolution":{"type":"boolean","value":"false"},'
        '"UseLosslessCompression":{"type":"boolean","value":"true"},'
        '"EmbedStandardFonts":{"type":"boolean","value":"true"},'
        '"MaxImageResolution":{"type":"long","value":"1200"}}'
    )

    cmd = [
        "soffice",
        "--headless",
        "--nologo",
        "--nolockcheck",
        "--nodefault",
        "--nofirststartwizard",
        "--norestore",
        f"-env:UserInstallation={profile_uri}",
        "--convert-to",
        convert_to,
        "--outdir",
        out_dir,
        pptx_path,
    ]

    proc = subprocess.run(cmd, capture_output=True, text=True, timeout=240)
    if proc.returncode != 0:
        raise RuntimeError(f"PDF conversion failed.\nSTDOUT:\n{proc.stdout}\nSTDERR:\n{proc.stderr}")

    pdf_name = os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf"
    pdf_path = os.path.join(out_dir, pdf_name)
    if not os.path.exists(pdf_path):
        raise RuntimeError("Conversion succeeded but PDF not found.")
    return pdf_path

# ----------------------------
# Streamlit UI
# ----------------------------
def main() -> None:
    st.set_page_config(page_title="Player Report Generator", layout="wide")
    st.title("Player Report Generator (PPTX / PDF)")


    # Session state
    st.session_state.setdefault("access_token", None)
    st.session_state.setdefault("api_base", None)
    st.session_state.setdefault("pptx_bytes", None)
    st.session_state.setdefault("pdf_bytes", None)
    st.session_state.setdefault("last_filename_base", None)

    # # Health checks
    # col_a, col_b, col_c = st.columns(3)
    # with col_a:
    #     st.caption("Template")
    #     st.write("✅" if os.path.exists(TEMPLATE_PPTX_PATH) else f"❌ Missing: {TEMPLATE_PPTX_PATH}")
    # with col_b:
    #     st.caption("bench.csv")
    #     st.write("✅" if os.path.exists(BENCH_CSV_PATH) else f"❌ Missing: {BENCH_CSV_PATH}")
    # with col_c:
    #     st.caption("PDF export")
    #     st.write("✅ LibreOffice found" if can_convert_to_pdf() else "⚠️ LibreOffice not found (PPTX only)")

    # Load bench.csv
    try:
        df_bench = load_bench_csv(BENCH_CSV_PATH)
    except Exception as e:
        st.error(str(e))
        st.stop()

    st.divider()

    # Step 1
    st.subheader("Step 1 — Generate access token (required)")
    cols = st.columns([1, 2])
    with cols[0]:
        gen = st.button("Generate access token", type="primary")
    with cols[1]:
        if st.session_state["access_token"]:
            st.success("Access token is set for this session.")
        else:
            st.warning(TOKEN_HINT)

    if gen:
        try:
            cfg = load_section_secrets()
            token = generate_access_token_from_secrets()
            st.session_state["access_token"] = token
            st.session_state["api_base"] = cfg.base_url
            st.success("Access token generated and stored for this session.")
        except Exception as e:
            st.session_state["access_token"] = None
            st.session_state["api_base"] = None
            st.error(f"Token generation failed: {e}")

    st.divider()

    # Step 2
    st.subheader("Step 2 — Select player and generate report")
    left, right = st.columns([1, 1])

    with left:
        player_label = st.selectbox(
            "Player (FC Den Bosch only)",
            options=[p["name"] for p in FC_DEN_BOSCH_PLAYERS],
            index=0,
        )
        player_id = next(p["player_id"] for p in FC_DEN_BOSCH_PLAYERS if p["name"] == player_label)

    with right:
        perf_file = st.file_uploader(
            "Optional performance chart (PNG/JPG). If not provided, it will be skipped.",
            type=["png", "jpg", "jpeg"],
            accept_multiple_files=False,
        )
        perf_bytes = perf_file.read() if perf_file else None
        if perf_file:
            st.success("Performance chart uploaded.")

    generate = st.button("Generate PPTX and PDF", type="primary")

    if generate:
        if not st.session_state["access_token"]:
            st.error("No valid access token. Generate the access token first.")
            st.stop()
        if not st.session_state.get("api_base"):
            st.error("Missing API base_url in secrets (key: base_url).")
            st.stop()

        token = st.session_state["access_token"]
        api_base = st.session_state["api_base"]

        with st.spinner("Generating report..."):
            try:
                with tempfile.TemporaryDirectory() as td:
                    base_name = re.sub(r"[^a-zA-Z0-9_-]+", "_", player_label).strip("_") or "player_report"
                    out_pptx_path = os.path.join(td, f"{base_name}.pptx")

                    meta = fill_template_full(
                        template_path=TEMPLATE_PPTX_PATH,
                        out_pptx_path=out_pptx_path,
                        df_bench=df_bench,
                        player_name_ui=player_label,
                        player_id=int(player_id),
                        api_base=api_base,
                        token=token,
                        performance_upload_bytes=perf_bytes,
                    )

                    with open(out_pptx_path, "rb") as f:
                        st.session_state["pptx_bytes"] = f.read()

                    st.session_state["last_filename_base"] = base_name

                    st.session_state["pdf_bytes"] = None
                    if can_convert_to_pdf():
                        try:
                            pdf_path = convert_pptx_to_pdf(out_pptx_path, td)
                            with open(pdf_path, "rb") as f:
                                st.session_state["pdf_bytes"] = f.read()
                        except Exception as e:
                            st.warning(f"PDF conversion failed (PPTX still available): {e}")

                st.success("Report generated.")
                with st.expander("Generation details"):
                    st.json(meta)

            except Exception as e:
                st.session_state["pptx_bytes"] = None
                st.session_state["pdf_bytes"] = None
                st.error(f"Generation failed: {e}")

    st.divider()

    # Downloads
    st.subheader("Downloads")
    base = st.session_state.get("last_filename_base") or "player_report"

    if st.session_state["pptx_bytes"]:
        st.download_button(
            "Download PPTX",
            data=st.session_state["pptx_bytes"],
            file_name=f"{base}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    else:
        st.info("Generate a report to enable PPTX download.")

    if st.session_state["pdf_bytes"]:
        st.download_button(
            "Download PDF",
            data=st.session_state["pdf_bytes"],
            file_name=f"{base}.pdf",
            mime="application/pdf",
        )


if __name__ == "__main__":
    main()
