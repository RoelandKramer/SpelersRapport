import os
from pptx import Presentation
from pptx.util import Inches

class PPTXProcessor:
    def __init__(self, template_path):
        self.template_path = template_path
        self.prs = Presentation(template_path)

    def replace_text(self, replacements: dict):
        """
        replaces {key} in the PPTX with value.
        Example: {'{NAME}': 'John Doe'}
        """
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for key, value in replacements.items():
                                if key in run.text:
                                    run.text = run.text.replace(key, value)

    def replace_image(self, placeholder_idx, image_path):
        """
        Replaces a picture placeholder by its index.
        """
        for slide in self.prs.slides:
            try:
                placeholder = slide.placeholders[placeholder_idx]
                if placeholder.placeholder_format.type == 18:  # 18 is PICTURE
                    placeholder.insert_picture(image_path)
            except (KeyError, AttributeError):
                continue

    def save(self, output_path):
        self.prs.save(output_path)
